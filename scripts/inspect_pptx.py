from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
import json

prs = Presentation(r'C:\Users\shrey\Downloads\InterviewLens_Presentation.pptx')
print(f'Slides: {len(prs.slides)}')
print(f'Size: {prs.slide_width.inches:.3f}" x {prs.slide_height.inches:.3f}"')
print(f'Size EMU: {prs.slide_width} x {prs.slide_height}')

for i, slide in enumerate(prs.slides):
    print(f'\n=== Slide {i+1}: layout={slide.slide_layout.name} ===')
    for sh in slide.shapes:
        name = getattr(sh, 'name', '')
        stype = sh.shape_type
        pos = f'({sh.left},{sh.top})'
        size = f'({sh.width},{sh.height})'
        fill_info = ''
        line_info = ''
        text_preview = ''
        try:
            fill = sh.fill
            if fill.type == 1:  # solid
                rgb = fill.fore_color.rgb
                fill_info = f' fill=#{rgb}'
        except: pass
        try:
            line = sh.line
            if line.color.type:
                rgb = line.color.rgb
                line_info = f' line=#{rgb}'
        except: pass
        try:
            if sh.has_text_frame:
                tf = sh.text_frame
                for para in tf.paragraphs[:2]:
                    for run in para.runs[:1]:
                        fc = ''
                        try: fc = f'#{run.font.color.rgb}'
                        except: pass
                        text_preview += f'[{run.font.size},{run.font.bold},{fc}:"{run.text[:40]}"] '
        except: pass
        print(f'  {name} type={stype} pos={pos} size={size}{fill_info}{line_info} {text_preview}')
