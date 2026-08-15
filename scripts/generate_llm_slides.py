"""
Generate 4 LLM & Reasoning slides matching the InterviewLens presentation style,
appended to the existing InterviewLens_Presentation.pptx.

Usage:
    python scripts/generate_llm_slides.py
Output: C:/Users/shrey/Downloads/InterviewLens_LLM_Slides.pptx  (4 new slides)
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Palette (from the existing deck) ──────────────────────────────────────────
BLUE       = RGBColor(0x2C, 0x7F, 0xB8)
DARK_NAVY  = RGBColor(0x14, 0x30, 0x4F)
MID_NAVY   = RGBColor(0x1A, 0x2B, 0x3C)
LIGHT_CARD = RGBColor(0xDC, 0xEB, 0xF7)   # #DCEBF7
ALT_CARD   = RGBColor(0xEA, 0xF2, 0xFA)   # #EAF2FA
GREEN_CARD = RGBColor(0xDD, 0xEA, 0xF6)   # #DDEAF6
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
MUTED      = RGBColor(0x5B, 0x7A, 0x99)
ORANGE     = RGBColor(0xF0, 0xB4, 0x29)   # accent for validation/warnings

# ── Slide dimensions (13.33" × 7.5") ──────────────────────────────────────────
W = 12192000   # EMU
H =  6858000   # EMU

# ── Helper shortcuts ───────────────────────────────────────────────────────────
def emu(inches): return int(inches * 914400)


def _set_para(tf, text, size_pt, bold=False, color: RGBColor = MID_NAVY,
              align=PP_ALIGN.LEFT, space_before=0):
    from pptx.util import Pt
    para = tf.paragraphs[0]
    para.alignment = align
    if space_before:
        para.space_before = Pt(space_before)
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return para


def _add_run(para, text, size_pt, bold=False, color: RGBColor = MID_NAVY):
    from pptx.util import Pt
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    return run


def _para(tf):
    """Append a new paragraph to a text frame."""
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsmap
    p = tf._txBody.add_p()
    return tf.paragraphs[-1]


def _add_shape(slide, l, t, w, h, fill: RGBColor | None = None,
               line: RGBColor | None = None, radius: int = 0) -> object:
    from pptx.util import Emu as E
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid() if fill else shape.fill.background()
    if fill: shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Emu(12700)
    else:
        shape.line.fill.background()
    return shape


def _add_textbox(slide, l, t, w, h) -> object:
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.line.fill.background()
    return txb


def _chrome(slide, section_title: str, page_num: int):
    """Add the standard chrome: blue circle icon + section bar + page number."""
    # Blue circle (top-left badge) — use an oval shape (MSO_SHAPE_TYPE = 9)
    from pptx.util import Emu as E
    circ = slide.shapes.add_shape(9, 514495, 419029, 615315, 615315)  # oval
    circ.fill.solid(); circ.fill.fore_color.rgb = BLUE
    circ.line.color.rgb = BLUE; circ.line.width = E(12700)

    # Section title bar (colored text, no background box needed)
    tb = _add_textbox(slide, 1252728, 588523, 9601200, 256032)
    tf = tb.text_frame
    tf.word_wrap = False
    _set_para(tf, section_title, 18, bold=True, color=BLUE)

    # Page number (bottom-right)
    pn = _add_textbox(slide, 11140135, 6528816, 548640, 256032)
    tf2 = pn.text_frame
    _set_para(tf2, str(page_num), 9, bold=False, color=MUTED, align=PP_ALIGN.RIGHT)


def _card(slide, l, t, w, h, fill: RGBColor = ALT_CARD):
    return _add_shape(slide, l, t, w, h, fill=fill, line=fill)


# ══════════════════════════════════════════════════════════════════════════════
# Slide A: "LLM & Reasoning — Architecture Overview"
# ══════════════════════════════════════════════════════════════════════════════
def slide_a(prs: Presentation, page: int):
    layout = prs.slide_layouts[0]   # blank
    slide  = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    _chrome(slide, "LLM & REASONING", page)

    # ── Headline ──────────────────────────────────────────────────────────────
    h = _add_textbox(slide, 1080412, 855054, 10037619, 476874)
    _set_para(h.text_frame, "From VLM to Evidence-Grounded LLM Reasoning", 25,
              bold=True, color=DARK_NAVY)

    # ── Sub-headline ───────────────────────────────────────────────────────────
    sub = _add_textbox(slide, 995567, 1380000, 10370934, 340000)
    _set_para(sub.text_frame,
              "The original Qwen2.5-VL multimodal VLM was replaced with a text-only "
              "Ollama-hosted LLM (Nemotron Mini 4B).  All reasoning is grounded "
              "in structured evidence — never free-form image captioning.",
              12, bold=False, color=MID_NAVY)

    # ── Pipeline flow: 5 cards in a row ───────────────────────────────────────
    labels = [
        ("Evidence\nPackage",    "Pose + audio +\nframing + background",  ALT_CARD),
        ("Prompt\nBuilder",      "Closed vocabulary,\ncoverage rules",     ALT_CARD),
        ("Nemotron\nMini 4B",    "Ollama local\nONNX-free inference",      LIGHT_CARD),
        ("Structured\nJSON Out", "observations /\nexplanations / suggestions", ALT_CARD),
        ("Validation\nLayer",    "Hallucination check\nreliability score", RGBColor(0xEA,0xF2,0xD6)),
    ]
    card_w = 1950000
    gap    = 190000
    total  = len(labels) * card_w + (len(labels) - 1) * gap
    start  = (W - total) // 2
    card_t = 2100000
    card_h = 1900000

    for i, (title, body, fill) in enumerate(labels):
        x = start + i * (card_w + gap)
        box = _card(slide, x, card_t, card_w, card_h, fill)
        # title
        tb = _add_textbox(slide, x + 80000, card_t + 130000, card_w - 160000, 500000)
        _set_para(tb.text_frame, title, 14, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)
        # body
        tb2 = _add_textbox(slide, x + 80000, card_t + 680000, card_w - 160000, 1100000)
        _set_para(tb2.text_frame, body, 10, bold=False, color=MID_NAVY, align=PP_ALIGN.CENTER)

        # Arrow (except after last)
        if i < len(labels) - 1:
            ax = x + card_w + 20000
            ay = card_t + card_h // 2 - 40000
            arr = _add_textbox(slide, ax, ay, gap - 40000, 120000)
            _set_para(arr.text_frame, "→", 18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    # ── Why text-only bullet ──────────────────────────────────────────────────
    why_box = _add_shape(slide, 502920, 4150000, W - 1005840, 690000,
                          fill=RGBColor(0xDC,0xEB,0xF7), line=None)
    tb3 = _add_textbox(slide, 660000, 4200000, W - 1320000, 620000)
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    _set_para(tf3,
              "▸  Text-only LLM removes torch / transformers / GPU dependencies entirely — "
              "runs on any laptop via Ollama.    "
              "▸  Evidence-grounded prompts force signal-type prefixes on every claim, "
              "preventing the hallucinations observed in VLM testing (0.25 reliability before fix).",
              11, bold=False, color=DARK_NAVY)


# ══════════════════════════════════════════════════════════════════════════════
# Slide B: "Evidence Assembly — What the LLM Sees"
# ══════════════════════════════════════════════════════════════════════════════
def slide_b(prs: Presentation, page: int):
    layout = prs.slide_layouts[0]
    slide  = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    _chrome(slide, "LLM & REASONING", page)

    h = _add_textbox(slide, 1080412, 855054, 10037619, 476874)
    _set_para(h.text_frame, "Evidence Assembly — Structured Input to the LLM", 25,
              bold=True, color=DARK_NAVY)

    # ── Two columns: left = inputs, right = EvidencePackage schema ────────────
    col_w = 5100000
    col_h = 3700000
    col_t = 1500000

    # Left: 4 signal-source boxes stacked
    sources = [
        ("📐  Pose signals",       "hands_near_face, arms_crossed,\nhead_tilt, body_lean, swaying …",    ALT_CARD),
        ("🖼️  Framing & BG",       "headroom_pct, centering_offset,\nbackground_distracting, low_light", ALT_CARD),
        ("🔊  Audio metrics",      "WPM, filler_word_count,\nlong_pause_count, low_mic_level",           ALT_CARD),
        ("📝  Transcript",         "Speaker text from ASR\n(real or synthetic placeholder)",             ALT_CARD),
    ]
    src_h = (col_h - 90000) // len(sources)
    for i, (title, body, fill) in enumerate(sources):
        y = col_t + i * (src_h + 30000)
        _card(slide, 502920, y, col_w - 100000, src_h - 10000, fill)
        tb = _add_textbox(slide, 620000, y + 60000, col_w - 340000, src_h - 120000)
        tf = tb.text_frame; tf.word_wrap = True
        p1 = tf.paragraphs[0]
        _add_run(p1, title + "   ", 11, bold=True, color=DARK_NAVY)
        _add_run(p1, body, 10, bold=False, color=MID_NAVY)

    # Arrow in middle
    arr = _add_textbox(slide, col_w + 350000, col_t + col_h // 2 - 150000, 400000, 300000)
    _set_para(arr.text_frame, "→", 28, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    # Right: EvidencePackage schema box
    ep_x = col_w + 800000
    ep_box = _card(slide, ep_x, col_t, col_w + 200000, col_h, fill=LIGHT_CARD)
    tb_ep = _add_textbox(slide, ep_x + 100000, col_t + 80000, col_w, col_h - 200000)
    tf_ep = tb_ep.text_frame; tf_ep.word_wrap = True
    ep_lines = [
        ("EvidencePackage", 14, True,  DARK_NAVY),
        ("", 8, False, MID_NAVY),
        ("  question:          str", 9,  False, MID_NAVY),
        ("  transcript:        Transcript", 9, False, MID_NAVY),
        ("  audio_metrics:     AudioMetrics", 9, False, MID_NAVY),
        ("  visual_events:     list[VisualEvent]", 9, False, MID_NAVY),
        ("  event_timestamps:  dict  ← signal summary, framing,", 9, False, MID_NAVY),
        ("                           background objects", 9, False, MID_NAVY),
        ("  selected_frames:   list[int]", 9, False, MID_NAVY),
        ("  frame_images:      list  ← unused by text-only LLM", 9, False, RGBColor(0xA0,0xB0,0xC0)),
    ]
    for i, (txt, size, bold, color) in enumerate(ep_lines):
        if i == 0:
            _set_para(tf_ep, txt, size, bold, color)
        else:
            p = _para(tf_ep)
            _add_run(p, txt, size, bold, color)

    # Explanation strip at bottom
    exp = _add_shape(slide, 502920, 5280000, W - 1005840, 600000,
                      fill=RGBColor(0xDC,0xEB,0xF7))
    tb_exp = _add_textbox(slide, 660000, 5320000, W - 1320000, 530000)
    _set_para(tb_exp.text_frame,
              "▸  assemble_evidence() and from_fused_evidence_json() both produce this contract.  "
              "▸  frame_images is populated but ignored — removing image tokens from the prompt "
              "eliminates appearance hallucinations (clothing, furniture) seen in VLM tests.",
              11, bold=False, color=DARK_NAVY)
    tb_exp.text_frame.word_wrap = True


# ══════════════════════════════════════════════════════════════════════════════
# Slide C: "Ollama + Nemotron Mini — Prompt Design"
# ══════════════════════════════════════════════════════════════════════════════
def slide_c(prs: Presentation, page: int):
    layout = prs.slide_layouts[0]
    slide  = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    _chrome(slide, "LLM & REASONING", page)

    h = _add_textbox(slide, 1080412, 855054, 10037619, 476874)
    _set_para(h.text_frame, "Ollama Nemotron Mini 4B — Prompt Engineering", 25,
              bold=True, color=DARK_NAVY)

    # ── Champion card (left) ───────────────────────────────────────────────────
    card_w = 3600000
    card_h = 3600000
    card_t = 1550000

    champ = _card(slide, 502920, card_t, card_w, card_h, fill=LIGHT_CARD)
    tb_c = _add_textbox(slide, 602920, card_t + 100000, card_w - 200000, 3400000)
    tf_c = tb_c.text_frame; tf_c.word_wrap = True
    _set_para(tf_c, "nemotron-mini  ✓ Champion", 13, bold=True, color=DARK_NAVY)
    lines_c = [
        "• NVIDIA Nemotron Mini 4B Instruct",
        "• Runs locally via Ollama (no API key)",
        "• ~2.7 GB one-time download",
        "• Instruction-tuned → follows strict JSON format",
        "• format='json' parameter enforces structured output",
        "• Falls back to mock if server offline",
    ]
    for line in lines_c:
        p = _para(tf_c)
        _add_run(p, line, 10, bold=False, color=MID_NAVY)

    # ── Challenger card (middle) ───────────────────────────────────────────────
    chal = _card(slide, 502920 + card_w + 200000, card_t, card_w, card_h,
                  fill=RGBColor(0xEE, 0xF1, 0xF4))
    tb_ch = _add_textbox(slide, 502920 + card_w + 300000, card_t + 100000,
                          card_w - 200000, 3400000)
    tf_ch = tb_ch.text_frame; tf_ch.word_wrap = True
    _set_para(tf_ch, "Qwen2.5-VL-3B  (replaced)", 13, bold=True,
              color=RGBColor(0x8B, 0x5C, 0xF6))
    lines_ch = [
        "• Multimodal VLM (vision + language)",
        "• Required torch / transformers / GPU",
        "• ~3B params — slow on CPU",
        "• Hallucinated clothing / furniture (0.25 reliability)",
        "• Small-model prompt-copying failure",
        "• Removed: no longer needed",
    ]
    for line in lines_ch:
        p = _para(tf_ch)
        _add_run(p, line, 10, bold=False, color=MID_NAVY)

    # ── Prompt design card (right) ─────────────────────────────────────────────
    right_x = 502920 + 2 * (card_w + 200000)
    right_w = W - right_x - 502920
    pcard = _card(slide, right_x, card_t, right_w, card_h, fill=ALT_CARD)
    tb_p = _add_textbox(slide, right_x + 100000, card_t + 100000,
                         right_w - 200000, 3400000)
    tf_p = tb_p.text_frame; tf_p.word_wrap = True
    _set_para(tf_p, "Prompt Design Rules", 13, bold=True, color=DARK_NAVY)
    rules = [
        "REQUIRED FORMAT: every observation",
        "must start with  <signal_type>: …",
        "",
        "REQUIRED COVERAGE: one observation",
        "per distinct signal type present",
        "",
        "Closed vocabulary — only tokens from",
        "SignalType enum are allowed",
        "",
        "Lighting / audio: concrete fixes",
        "mandated (not vague comments)",
    ]
    for rule in rules:
        p = _para(tf_p)
        _add_run(p, rule, 9, bold=False, color=MID_NAVY if rule else MUTED)

    # ── Bottom: champion label ─────────────────────────────────────────────────
    chl = _add_textbox(slide, 502920 + 800000, card_t - 380000, 1280262, 315469)
    _set_para(chl.text_frame, "Champion", 16, bold=True, color=DARK_NAVY)
    cha = _add_textbox(slide, 502920 + card_w + 200000 + 800000, card_t - 380000, 1280262, 315469)
    _set_para(cha.text_frame, "Challenger", 16, bold=True, color=DARK_NAVY)

    # Bottom explanation strip
    exp = _add_shape(slide, 502920, 5280000, W - 1005840, 600000,
                      fill=RGBColor(0xDC,0xEB,0xF7))
    tb_e = _add_textbox(slide, 660000, 5320000, W - 1320000, 530000)
    _set_para(tb_e.text_frame,
              "▸  REQUIRED FORMAT forces signal-type prefixes — a static example in context "
              "caused the small VLM to copy the example verbatim (known small-model failure mode).  "
              "▸  Closed vocabulary ties every claim back to enum values validated downstream.",
              11, bold=False, color=DARK_NAVY)
    tb_e.text_frame.word_wrap = True


# ══════════════════════════════════════════════════════════════════════════════
# Slide D: "Evidence Validation & Coaching Report"
# ══════════════════════════════════════════════════════════════════════════════
def slide_d(prs: Presentation, page: int):
    layout = prs.slide_layouts[0]
    slide  = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    _chrome(slide, "LLM & REASONING", page)

    h = _add_textbox(slide, 1080412, 855054, 10037619, 476874)
    _set_para(h.text_frame, "Evidence Validation & Coaching Report Generation", 25,
              bold=True, color=DARK_NAVY)

    # ── Left: validation checks stack ─────────────────────────────────────────
    checks = [
        ("#1  Allowed-category check",
         "Every claim must contain a word from ALLOWED_KEYWORDS\n"
         "(all SignalType enum values + 60 plain-English terms).", ALT_CARD),
        ("#2  Timestamp validity check",
         "Any timestamp referenced in a claim must fall within\n"
         "0 – clip_end seconds.", ALT_CARD),
        ("#3  Confidence gate",
         "Low-confidence detections (< 0.5) lower reliability;\n"
         "high-confidence events must be present.", ALT_CARD),
        ("#4  Sufficient-evidence check",
         "Strong claims require matching evidence data points.\n"
         "Reliability score: max(0, 1.0 − 0.15 × failures).",
         RGBColor(0xEA,0xF2,0xD6)),
    ]
    chk_w = 5200000
    chk_h = 1100000
    for i, (title, body, fill) in enumerate(checks):
        y = 1550000 + i * (chk_h + 80000)
        _card(slide, 502920, y, chk_w, chk_h, fill)
        tb = _add_textbox(slide, 640000, y + 80000, chk_w - 280000, chk_h - 160000)
        tf = tb.text_frame; tf.word_wrap = True
        _set_para(tf, title, 11, bold=True, color=DARK_NAVY)
        p2 = _para(tf)
        _add_run(p2, body, 10, bold=False, color=MID_NAVY)

    # ── Right: CoachingReport schema ───────────────────────────────────────────
    rpt_x = chk_w + 800000
    rpt_w = W - rpt_x - 400000
    rpt_h = 4 * (chk_h + 80000) - 80000
    _card(slide, rpt_x, 1550000, rpt_w, rpt_h, fill=LIGHT_CARD)
    tb_r = _add_textbox(slide, rpt_x + 100000, 1650000, rpt_w - 200000, rpt_h - 200000)
    tf_r = tb_r.text_frame; tf_r.word_wrap = True
    _set_para(tf_r, "CoachingReport", 13, bold=True, color=DARK_NAVY)
    report_fields = [
        ("", False),
        ("  duration_s              float", False),
        ("  valid_tracking_pct      float  (0–100%)", False),
        ("  speaking_rate_wpm       float", False),
        ("  filler_word_count       int", False),
        ("  long_pause_count        int", False),
        ("  detected_signals        int", False),
        ("", False),
        ("  strengths               list[str]  ← observations", False),
        ("  improvements            list[str]  ← suggestions", False),
        ("  timeline                list[{time_s, label}]", False),
        ("  reliability_score       float  (0.0–1.0)", False),
        ("", False),
        ("  → PDF download  ✓", True),
        ("  → Streamlit dashboard  ✓", True),
    ]
    for txt, bold in report_fields:
        p = _para(tf_r)
        _add_run(p, txt, 9, bold=bold,
                  color=BLUE if bold else MID_NAVY)

    # ── Bottom strip ───────────────────────────────────────────────────────────
    exp = _add_shape(slide, 502920, 5280000, W - 1005840, 600000,
                      fill=RGBColor(0xDC,0xEB,0xF7))
    tb_e = _add_textbox(slide, 660000, 5320000, W - 1320000, 530000)
    _set_para(tb_e.text_frame,
              "▸  Reliability ≥ 0.75 = High confidence (green).  0.50–0.75 = Moderate (amber).  "
              "< 0.50 = Low (red).  Report shows coaching content at ≥ 0.3; below that a fallback "
              "message is shown.  "
              "▸  Full coaching report exported as a branded PDF from the Streamlit dashboard.",
              11, bold=False, color=DARK_NAVY)
    tb_e.text_frame.word_wrap = True


# ══════════════════════════════════════════════════════════════════════════════
# Main
# ══════════════════════════════════════════════════════════════════════════════
def main():
    src  = Path(r"C:\Users\shrey\Downloads\InterviewLens_Presentation.pptx")
    dest = Path(r"C:\Users\shrey\Downloads\InterviewLens_LLM_Slides.pptx")

    prs = Presentation(str(src))
    n   = len(prs.slides)
    print(f"Loaded {src.name} — {n} existing slides")

    slide_a(prs, n + 1)
    print(f"  Added slide {n+1}: LLM Architecture Overview")
    slide_b(prs, n + 2)
    print(f"  Added slide {n+2}: Evidence Assembly")
    slide_c(prs, n + 3)
    print(f"  Added slide {n+3}: Ollama + Prompt Design")
    slide_d(prs, n + 4)
    print(f"  Added slide {n+4}: Validation & Coaching Report")

    prs.save(str(dest))
    size_mb = dest.stat().st_size / 1024 / 1024
    print(f"\nSaved → {dest}  ({size_mb:.1f} MB, {n+4} slides total)")


if __name__ == "__main__":
    main()
